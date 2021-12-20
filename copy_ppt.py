from pptx import Presentation
import copy
import os


def copy_slide(des, src_slide, placeholdertexts):
    """
    从源slide copy并根据placeholdertext赋值
    """
    print(src_slide.slide_layout)
    curr_slide = des.slides.add_slide(src_slide.slide_layout)  # 直接使用src的layout
    # remove placeholder
    for placeholder in curr_slide.placeholders:
        sp = placeholder._sp
        sp.getparent().remove(sp)

    for shp in src_slide.shapes:
        copy_shape(shp, curr_slide.shapes)

    if placeholdertexts is not None:  # 定义文件中有设置placeholder的值
        index = 0
        for placeholder in curr_slide.placeholders:
            if index < len(placeholdertexts):  # 超出不再赋值
                placeholder.text = placeholdertexts[index]
                index = index+1


def copy_shape(shp, parent_shapes):
    if '图片' in shp.name or 'Picture' in shp.name:
        # save image
        add_image(shp, parent_shapes)
    else:
        el = shp.element
        newel = copy.deepcopy(el)
        parent_shapes._spTree.insert_element_before(newel, 'p:extLst')


def add_image(shp, parent_shapes):
    image_path = shp.name+'.jpg'
    with open(image_path, 'wb') as f:
        f.write(shp.image.blob)

    # add image to dict
    parent_shapes.add_picture(
        image_path, shp.left, shp.top, shp.width, shp.height)
    os.remove(image_path)


def copy_ppt_index(des, src, index, placeholdertexts):
    """
    从指定ppt copy其中一页，从1开始，所以要-1
    """
    src_slide = src.slides[index-1]
    copy_slide(des, src_slide, placeholdertexts)


def copy_ppt(des, src, placeholdertexts):
    """
    直接copy整个ppt
    """
    for slide in src.slides:
        copy_slide(des, slide, placeholdertexts)
